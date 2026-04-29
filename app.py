"""
app.py — BA Agent Chat Interface
----------------------------------
A browser-based chat UI built with Streamlit.
The BA can upload documents and ask questions
without touching the terminal.

Run with:
    streamlit run app.py
"""

import os
import streamlit as st
import chromadb
from chromadb.utils import embedding_functions
import anthropic
import fitz
from docx import Document
from pptx import Presentation
import tempfile

# ─────────────────────────────────────────
# SECRETS — works on Render, Streamlit Cloud, and locally
# ─────────────────────────────────────────
try:
    from dotenv import load_dotenv
    load_dotenv()
except Exception:
    pass

def _get_secret(key):
    val = os.environ.get(key, "")
    if val:
        return val
    try:
        return st.secrets.get(key, "")
    except Exception:
        return ""

ANTHROPIC_API_KEY = _get_secret("ANTHROPIC_API_KEY")
TAVILY_API_KEY    = _get_secret("TAVILY_API_KEY")

# ─────────────────────────────────────────
# SETTINGS
# ─────────────────────────────────────────
CHROMA_FOLDER   = "chroma_db"
COLLECTION_NAME = "ba_documents"
UPLOADS_FOLDER  = "uploads"
CHUNK_SIZE      = 400
CHUNK_OVERLAP   = 50
TOP_K_RESULTS   = 5

# ─────────────────────────────────────────
# PAGE CONFIG
# ─────────────────────────────────────────
st.set_page_config(
    page_title = "BA Document Agent",
    page_icon  = "📋",
    layout     = "wide"
)

# ─────────────────────────────────────────
# CUSTOM STYLING
# ─────────────────────────────────────────
st.markdown("""
<style>
    .main-header {
        background: linear-gradient(135deg, #1e3a5f 0%, #2d6a9f 100%);
        padding: 20px 28px;
        border-radius: 12px;
        margin-bottom: 24px;
        color: white;
    }
    .main-header h1 { margin: 0; font-size: 24px; font-weight: 600; }
    .main-header p  { margin: 4px 0 0; font-size: 13px; opacity: 0.85; }

    .source-badge {
        background: #f0f4ff;
        border: 1px solid #c7d7f5;
        border-radius: 6px;
        padding: 4px 10px;
        font-size: 12px;
        color: #2d5faa;
        display: inline-block;
        margin: 3px 3px 3px 0;
    }
    .web-badge {
        background: #f0fff4;
        border: 1px solid #b7e8c7;
        border-radius: 6px;
        padding: 4px 10px;
        font-size: 12px;
        color: #1e7e3a;
        display: inline-block;
        margin: 3px 3px 3px 0;
    }
    .status-box {
        background: #f8faff;
        border-left: 3px solid #2d6a9f;
        padding: 10px 14px;
        border-radius: 0 8px 8px 0;
        font-size: 13px;
        margin-bottom: 8px;
    }
    .chunk-count {
        background: #e8f4fd;
        border-radius: 20px;
        padding: 2px 10px;
        font-size: 12px;
        color: #1a5f8a;
        font-weight: 500;
    }
    div[data-testid="stChatMessage"] {
        border-radius: 12px;
        margin-bottom: 8px;
    }
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────
# SESSION STATE — persists across reruns
# ─────────────────────────────────────────
if "messages"           not in st.session_state: st.session_state.messages           = []
if "collection"         not in st.session_state: st.session_state.collection         = None
if "ingested_files"     not in st.session_state: st.session_state.ingested_files     = []
if "total_chunks"       not in st.session_state: st.session_state.total_chunks       = 0

# ─────────────────────────────────────────
# HELPER — Parsing functions (from ingest.py)
# ─────────────────────────────────────────

def extract_pdf(file_bytes, filename):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    blocks = []
    for i, page in enumerate(doc, 1):
        text = page.get_text().strip()
        if text:
            blocks.append({"text": text, "source": filename, "page": f"Page {i}"})
    return blocks

def extract_docx(file_bytes, filename):
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    doc = Document(tmp_path)
    os.unlink(tmp_path)
    blocks, buffer, para_num = [], "", 0
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text: continue
        para_num += 1
        buffer += " " + text
        if para_num % 5 == 0:
            blocks.append({"text": buffer.strip(), "source": filename, "page": f"Para {para_num-4}–{para_num}"})
            buffer = ""
    if buffer.strip():
        blocks.append({"text": buffer.strip(), "source": filename, "page": f"Para {para_num}"})
    return blocks

def extract_pptx(file_bytes, filename):
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    prs = Presentation(tmp_path)
    os.unlink(tmp_path)
    blocks = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [p.text.strip() for shape in slide.shapes if shape.has_text_frame for p in shape.text_frame.paragraphs if p.text.strip()]
        if texts:
            blocks.append({"text": " ".join(texts), "source": filename, "page": f"Slide {i}"})
    return blocks

def split_chunks(block):
    words  = block["text"].split()
    chunks = []
    start  = 0
    while start < len(words):
        chunk_words = words[start:start + CHUNK_SIZE]
        chunks.append({"text": " ".join(chunk_words), "source": block["source"], "location": block["page"]})
        start += CHUNK_SIZE - CHUNK_OVERLAP
    return chunks

# ─────────────────────────────────────────
# HELPER — Ingest uploaded files
# ─────────────────────────────────────────

def ingest_files(uploaded_files):
    embedding_fn = embedding_functions.DefaultEmbeddingFunction()
    client       = chromadb.PersistentClient(path=CHROMA_FOLDER)

    try:    client.delete_collection(name=COLLECTION_NAME)
    except: pass

    collection = client.create_collection(name=COLLECTION_NAME, embedding_function=embedding_fn)

    all_chunks    = []
    file_names    = []
    progress_bar  = st.progress(0)
    status_text   = st.empty()

    for idx, uploaded_file in enumerate(uploaded_files):
        file_bytes = uploaded_file.read()
        filename   = uploaded_file.name
        ext        = filename.lower().split(".")[-1]
        file_names.append(filename)

        status_text.markdown(f"**Reading** `{filename}`...")

        try:
            if   ext == "pdf":  blocks = extract_pdf(file_bytes, filename)
            elif ext == "docx": blocks = extract_docx(file_bytes, filename)
            elif ext == "pptx": blocks = extract_pptx(file_bytes, filename)
            else: continue

            for block in blocks:
                all_chunks.extend(split_chunks(block))

        except Exception as e:
            st.error(f"❌ Error reading {filename}: {e}")

        progress_bar.progress((idx + 1) / len(uploaded_files))

    status_text.markdown(f"**Saving** {len(all_chunks)} chunks to memory...")

    batch_size = 100
    for i in range(0, len(all_chunks), batch_size):
        batch = all_chunks[i:i+batch_size]
        collection.add(
            documents = [c["text"]     for c in batch],
            metadatas = [{"source": c["source"], "location": c["location"]} for c in batch],
            ids       = [f"chunk_{i+j}" for j in range(len(batch))]
        )

    progress_bar.empty()
    status_text.empty()

    return collection, file_names, len(all_chunks)

# ─────────────────────────────────────────
# HELPER — Search & Answer
# ─────────────────────────────────────────

def search_documents(collection, question):
    results = collection.query(query_texts=[question], n_results=TOP_K_RESULTS)
    return [{"text": results["documents"][0][i], "source": results["metadatas"][0][i].get("source",""), "location": results["metadatas"][0][i].get("location","")} for i in range(len(results["documents"][0]))]

def search_web(question):
    if not TAVILY_API_KEY: return []
    try:
        from tavily import TavilyClient
        results = TavilyClient(api_key=TAVILY_API_KEY).search(query=question, max_results=3)
        return [{"text": r.get("content",""), "source": r.get("url","Web"), "location": "Web"} for r in results.get("results",[])]
    except:
        return []

def get_answer(question, doc_chunks, web_chunks):
    context_parts = []
    if doc_chunks:
        context_parts.append("=== FROM YOUR UPLOADED DOCUMENTS ===")
        for i, c in enumerate(doc_chunks, 1):
            context_parts.append(f"[{i}] Source: {c['source']} | {c['location']}\n{c['text']}")
    if web_chunks:
        context_parts.append("\n=== FROM WEB SEARCH ===")
        for i, c in enumerate(web_chunks, len(doc_chunks)+1):
            context_parts.append(f"[{i}] Source: {c['source']}\n{c['text']}")
    context = "\n\n".join(context_parts)

    history = [{"role": m["role"], "content": m["content"]} for m in st.session_state.messages[-10:]]
    history.append({"role": "user", "content": f"Context:\n{context}\n\nQuestion: {question}\n\nAnswer with citations."})

 
    client = anthropic.Anthropic(api_key="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")
    resp   = client.messages.create(
        model   = "claude-sonnet-4-5",
        max_tokens = 1500,
        system  = """You are an expert Business Analyst assistant. Help BAs understand their project documents.
RULES:
1. Answer based on the provided context (uploaded documents) first.
2. Cite sources clearly using format: 📄 [filename | location]
3. Be clear, structured, and professional.
4. Flag conflicting information across documents.
5. Suggest follow-up questions at the end when relevant.
6. If context is insufficient, say so and suggest what document might help.""",
        messages = history
    )
    return resp.content[0].text

# ─────────────────────────────────────────
# SIDEBAR — File upload + status
# ─────────────────────────────────────────

with st.sidebar:
    st.markdown("### 📂 Upload Documents")
    st.markdown("Supported: **PDF, DOCX, PPTX**")

    uploaded_files = st.file_uploader(
        label      = "Drop files here",
        type       = ["pdf", "docx", "pptx"],
        accept_multiple_files = True,
        label_visibility = "collapsed"
    )

    if uploaded_files:
        if st.button("📥 Process Documents", use_container_width=True, type="primary"):
            with st.spinner("Processing your documents..."):
                collection, names, total = ingest_files(uploaded_files)
                st.session_state.collection     = collection
                st.session_state.ingested_files = names
                st.session_state.total_chunks   = total
                st.session_state.messages       = []  # reset chat on new docs
            st.success(f"✅ Ready! {total} chunks indexed.")

    st.divider()

    # Show loaded documents
    if st.session_state.ingested_files:
        st.markdown("### 📚 Loaded Documents")
        st.markdown(f'<span class="chunk-count">{st.session_state.total_chunks} chunks in memory</span>', unsafe_allow_html=True)
        st.markdown("")
        for f in st.session_state.ingested_files:
            ext = f.split(".")[-1].upper()
            color = {"PDF":"🔴","DOCX":"🔵","PPTX":"🟠"}.get(ext,"📄")
            st.markdown(f"{color} `{f}`")

    st.divider()

    # Web search toggle
    use_web = st.toggle(
        "🌐 Web Search",
        value = bool(TAVILY_API_KEY),
        help  = "Search the web for additional context (requires Tavily API key)"
    )

    st.divider()

    # Clear chat button
    if st.button("🗑️ Clear Chat", use_container_width=True):
        st.session_state.messages = []
        st.rerun()

    st.markdown("")
    st.markdown('<div style="font-size:11px;color:#999;">BA Document Agent v1.0<br>Powered by Claude + ChromaDB</div>', unsafe_allow_html=True)

# ─────────────────────────────────────────
# MAIN AREA — Header + Chat
# ─────────────────────────────────────────

st.markdown("""
<div class="main-header">
    <h1>📋 BA Document Agent</h1>
    <p>Upload your project documents and ask anything — get cited answers instantly</p>
</div>
""", unsafe_allow_html=True)

# Show welcome message if no documents loaded
if not st.session_state.collection:
    st.info("👈 Upload your PDF, Word, or PowerPoint files in the sidebar to get started.")

    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown("""
        **📄 What you can upload**
        - Requirements docs
        - Meeting notes
        - Project charters
        - Stakeholder decks
        - RFPs & proposals
        """)
    with col2:
        st.markdown("""
        **💬 What you can ask**
        - Summarise this document
        - What are the project risks?
        - What decisions were made?
        - Compare requirements across docs
        - What are the open action items?
        """)
    with col3:
        st.markdown("""
        **✅ What you get back**
        - Clear, structured answers
        - Exact source citations
        - Page / slide references
        - Web search context
        - Follow-up suggestions
        """)

# ─────────────────────────────────────────
# CHAT HISTORY — Display past messages
# ─────────────────────────────────────────

for message in st.session_state.messages:
    with st.chat_message(message["role"], avatar="🧑‍💼" if message["role"]=="user" else "🤖"):
        st.markdown(message["content"])

        # Show source badges under agent messages
        if message["role"] == "assistant" and "sources" in message:
            st.markdown("")
            sources_html = ""
            for s in message["sources"]:
                if s["location"] == "Web":
                    sources_html += f'<span class="web-badge">🌐 {s["source"][:50]}</span>'
                else:
                    sources_html += f'<span class="source-badge">📄 {s["source"]} — {s["location"]}</span>'
            st.markdown(sources_html, unsafe_allow_html=True)

# ─────────────────────────────────────────
# CHAT INPUT
# ─────────────────────────────────────────

if prompt := st.chat_input("Ask anything about your documents..."):

    if not st.session_state.collection:
        st.warning("⚠️ Please upload and process documents first using the sidebar.")
        st.stop()

    # Show user message
    with st.chat_message("user", avatar="🧑‍💼"):
        st.markdown(prompt)
    st.session_state.messages.append({"role": "user", "content": prompt})

    # Generate answer
    with st.chat_message("assistant", avatar="🤖"):
        with st.spinner("Searching documents and generating answer..."):

            doc_chunks = search_documents(st.session_state.collection, prompt)
            web_chunks = search_web(prompt) if use_web else []
            answer     = get_answer(prompt, doc_chunks, web_chunks)

        st.markdown(answer)

        # Show source badges
        all_sources = doc_chunks + web_chunks
        if all_sources:
            st.markdown("")
            sources_html = ""
            for s in all_sources:
                if s["location"] == "Web":
                    sources_html += f'<span class="web-badge">🌐 {s["source"][:60]}</span>'
                else:
                    sources_html += f'<span class="source-badge">📄 {s["source"]} — {s["location"]}</span>'
            st.markdown(sources_html, unsafe_allow_html=True)

    # Save assistant message with sources
    st.session_state.messages.append({
        "role"   : "assistant",
        "content": answer,
        "sources": all_sources
    })
