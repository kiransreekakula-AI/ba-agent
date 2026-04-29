"""
ingest.py — BA Agent Document Parser
--------------------------------------
Reads PDF, Word (.docx), and PowerPoint (.pptx) files,
breaks them into chunks, and stores them in ChromaDB
so the agent can search them later.

Usage:
    python3 ingest.py
"""

import os
import fitz                          # PyMuPDF — reads PDFs
from docx import Document            # reads Word files
from pptx import Presentation        # reads PowerPoint files
from dotenv import load_dotenv
import chromadb                      # the searchable memory database
from chromadb.utils import embedding_functions

load_dotenv()  # reads your .env file

# ─────────────────────────────────────────
# SETTINGS — adjust these if needed
# ─────────────────────────────────────────
UPLOADS_FOLDER = "uploads"           # folder where BA drops files
CHROMA_FOLDER  = "chroma_db"        # folder where memory is stored
CHUNK_SIZE     = 400                 # words per chunk (sweet spot for BA docs)
CHUNK_OVERLAP  = 50                  # words overlap between chunks
COLLECTION_NAME = "ba_documents"    # name of the memory collection

# ─────────────────────────────────────────
# STEP 1 — Extract text from each file type
# ─────────────────────────────────────────

def extract_from_pdf(filepath):
    """Read a PDF and return list of {text, source, page} dicts."""
    chunks_raw = []
    doc = fitz.open(filepath)
    filename = os.path.basename(filepath)
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text().strip()
        if text:  # skip blank pages
            chunks_raw.append({
                "text"  : text,
                "source": filename,
                "page"  : f"Page {page_num}"
            })
    print(f"  ✅ PDF: {filename} — {len(doc)} pages extracted")
    return chunks_raw


def extract_from_docx(filepath):
    """Read a Word doc and return list of {text, source, page} dicts."""
    chunks_raw = []
    doc = Document(filepath)
    filename = os.path.basename(filepath)
    para_num = 0
    buffer = ""

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        para_num += 1
        buffer += " " + text

        # group every 5 paragraphs into one raw chunk
        if para_num % 5 == 0:
            chunks_raw.append({
                "text"  : buffer.strip(),
                "source": filename,
                "page"  : f"Para {para_num - 4}–{para_num}"
            })
            buffer = ""

    # catch any remaining text
    if buffer.strip():
        chunks_raw.append({
            "text"  : buffer.strip(),
            "source": filename,
            "page"  : f"Para {para_num}"
        })

    print(f"  ✅ DOCX: {filename} — {para_num} paragraphs extracted")
    return chunks_raw


def extract_from_pptx(filepath):
    """Read a PowerPoint and return list of {text, source, page} dicts."""
    chunks_raw = []
    prs = Presentation(filepath)
    filename = os.path.basename(filepath)

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            chunks_raw.append({
                "text"  : " ".join(texts),
                "source": filename,
                "page"  : f"Slide {slide_num}"
            })

    print(f"  ✅ PPTX: {filename} — {len(prs.slides)} slides extracted")
    return chunks_raw


# ─────────────────────────────────────────
# STEP 2 — Split text into smaller chunks
# ─────────────────────────────────────────

def split_into_chunks(raw_item):
    """
    Take one raw extracted block and split into overlapping chunks.
    Each chunk remembers where it came from (source + page).
    """
    words = raw_item["text"].split()
    chunks = []
    start = 0

    while start < len(words):
        end = start + CHUNK_SIZE
        chunk_words = words[start:end]
        chunk_text  = " ".join(chunk_words)

        chunks.append({
            "text"    : chunk_text,
            "source"  : raw_item["source"],
            "location": raw_item["page"]
        })

        # move forward, but overlap by CHUNK_OVERLAP words
        start += CHUNK_SIZE - CHUNK_OVERLAP

    return chunks


# ─────────────────────────────────────────
# STEP 3 — Store chunks in ChromaDB
# ─────────────────────────────────────────

def store_in_chromadb(all_chunks):
    """Save all chunks into ChromaDB so the agent can search them."""

    # Use a simple local embedding model (no extra API key needed)
    embedding_fn = embedding_functions.DefaultEmbeddingFunction()

    client     = chromadb.PersistentClient(path=CHROMA_FOLDER)

    # Delete old collection if it exists (fresh start on each ingest)
    try:
        client.delete_collection(name=COLLECTION_NAME)
        print(f"\n  🗑️  Old memory cleared — starting fresh")
    except:
        pass

    collection = client.create_collection(
        name               = COLLECTION_NAME,
        embedding_function = embedding_fn
    )

    # Prepare data for ChromaDB
    documents  = [c["text"]     for c in all_chunks]
    metadatas  = [{"source": c["source"], "location": c["location"]} for c in all_chunks]
    ids        = [f"chunk_{i}"  for i in range(len(all_chunks))]

    # Store in batches of 100 to avoid memory issues
    batch_size = 100
    for i in range(0, len(all_chunks), batch_size):
        collection.add(
            documents = documents[i:i+batch_size],
            metadatas = metadatas[i:i+batch_size],
            ids       = ids[i:i+batch_size]
        )

    print(f"  💾 Stored {len(all_chunks)} chunks in memory")
    return collection


# ─────────────────────────────────────────
# MAIN — Run everything
# ─────────────────────────────────────────

def ingest_all():
    print("\n" + "="*50)
    print("  BA Agent — Document Ingestion")
    print("="*50)

    # Check uploads folder exists
    if not os.path.exists(UPLOADS_FOLDER):
        os.makedirs(UPLOADS_FOLDER)
        print(f"\n  ⚠️  Created '{UPLOADS_FOLDER}/' folder.")
        print(f"  👉 Drop your PDF, DOCX, PPTX files into '{UPLOADS_FOLDER}/' and run again.\n")
        return

    # Collect all supported files
    supported = (".pdf", ".docx", ".pptx")
    files = [f for f in os.listdir(UPLOADS_FOLDER) if f.lower().endswith(supported)]

    if not files:
        print(f"\n  ⚠️  No files found in '{UPLOADS_FOLDER}/'")
        print(f"  👉 Drop your PDF, DOCX, or PPTX files there and run again.\n")
        return

    print(f"\n  📂 Found {len(files)} file(s) to process:\n")

    all_chunks = []

    for filename in files:
        filepath = os.path.join(UPLOADS_FOLDER, filename)
        ext = filename.lower().split(".")[-1]

        try:
            if ext == "pdf":
                raw_blocks = extract_from_pdf(filepath)
            elif ext == "docx":
                raw_blocks = extract_from_docx(filepath)
            elif ext == "pptx":
                raw_blocks = extract_from_pptx(filepath)
            else:
                continue

            # Split each block into chunks
            for block in raw_blocks:
                all_chunks.extend(split_into_chunks(block))

        except Exception as e:
            print(f"  ❌ Error reading {filename}: {e}")

    print(f"\n  📄 Total chunks created: {len(all_chunks)}")

    # Store everything in ChromaDB
    print("\n  💾 Saving to memory (ChromaDB)...")
    store_in_chromadb(all_chunks)

    print("\n" + "="*50)
    print("  ✅ Ingestion complete! Agent memory is ready.")
    print("="*50 + "\n")


if __name__ == "__main__":
    ingest_all()
